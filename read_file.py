from pdfminer.high_level import extract_text
from personal_info import info_search
from win32com.client import Dispatch
from pptx import Presentation
import pandas as pd
import docx
import csv
import os
import sys

sys.tracebacklimit = 0


def read_pdf(pdf_path):
    try:
        tmp_path = './tmp/pdf_tmp.txt'
        text = extract_text(pdf_path)
        tf = open(tmp_path, 'w', encoding='UTF-8')
        tf.write(text)
        tf.close()
        tf = open(tmp_path, 'r', encoding='UTF-8')
        rusult = info_search(tf.read())
        # print(len(test['username']))
        tf.close()
    except Exception as e:
        print(e)
        rusult = {"username": [], "id": [], "ph_no": [],
                  "em": [], "address": [], "bir": []}
        pass
    try:
        # 刪除暫存檔
        # os.remove(pdf_path)
        os.remove(tmp_path)
        os.remove(pdf_path)
    except:
        pass
    return rusult


def read_docx(docx_path):
    try:
        tmp_path = './tmp/doc_tmp.txt'
        file = docx.Document(docx_path)
        tf = open(tmp_path, 'w', encoding='UTF-8')
        for para in file.paragraphs:  # 獲取一般內容
            tf.write(para.text + '\n')
        tables = file.tables   # 獲取表格內容
        for table in tables:
            for row in table.rows:  # 遍歷表格的所有行
                for cell in row.cells:
                    tf.write(cell.text + '\n')
        file.save(docx_path)
        tf.close()
        tf = open(tmp_path, 'r', encoding='UTF-8')
        rusult = info_search(tf.read())
        tf.close()
    except Exception as e:
        print(e)
        rusult = {"username": [], "id": [], "ph_no": [],
                  "em": [], "address": [], "bir": []}
        pass
    try:
        # 刪除暫存檔
        os.remove(tmp_path)
        os.remove(docx_path)
    except:
        pass
    return rusult


def read_rtf(rtf_path):
    try:
        tmp_path = './tmp/rtf_tmp.txt'
        tf = open(tmp_path, 'w', encoding='UTF-8')
        word = Dispatch("Word.Application")
        wordfile = word.Documents.Open(sys.path[0] + rtf_path)
        for para in wordfile.paragraphs:
            tf.write(para.Range.Text + '\n')
        wordfile.Close()
        word.Quit()
        tf.close()
        tf = open(tmp_path, 'r', encoding='UTF-8')
        rusult = info_search(tf.read())
        tf.close()
    except Exception as e:
        print(e)
        rusult = {"username": [], "id": [], "ph_no": [],
                  "em": [], "address": [], "bir": []}
        pass
    try:
        # 刪除暫存檔
        os.remove(tmp_path)
        os.remove(rtf_path)
    except:
        pass
    return rusult


def read_xlsx(xlsx_path):
    try:
        tmp_path = './tmp/xlsx_tmp.txt'
        text = pd.read_excel(xlsx_path, sheet_name=None)
        tf = open(tmp_path, 'w', encoding='UTF-8')
        tf.write(str(text))
        tf.close()
        tf = open(tmp_path, 'r', encoding='UTF-8')
        rusult = info_search(tf.read())
        # print(len(test['username']))
        tf.close()
    except Exception as e:
        print(e)
        rusult = {"username": [], "id": [], "ph_no": [],
                  "em": [], "address": [], "bir": []}
        pass
    try:
        # 刪除暫存檔
        # os.remove(pdf_path)
        os.remove(tmp_path)
        os.remove(xlsx_path)
    except:
        pass
    return rusult


def read_csv(csv_path):
    try:
        tmp_path = './tmp/csv_tmp.txt'
        tf = open(tmp_path, 'w', encoding='UTF-8')
        with open(csv_path, newline='', encoding='UTF-8') as csvfile:
            rows = csv.reader(csvfile)
            for row in rows:
                tf.write(str(row))
        tf.close()
        tf = open(tmp_path, 'r', encoding='UTF-8')
        rusult = info_search(tf.read())
        # print(len(test['username']))
        tf.close()
    except Exception as e:
        print(e)
        rusult = {"username": [], "id": [], "ph_no": [],
                  "em": [], "address": [], "bir": []}
        pass
    try:
        # 刪除暫存檔
        # os.remove(pdf_path)
        os.remove(tmp_path)
        os.remove(csv_path)
    except:
        pass
    return rusult


def read_pptx(pptx_path):
    try:
        tmp_path = './tmp/pptx_tmp.txt'
        tf = open(tmp_path, 'w', encoding='UTF-8')
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    tf.write(shape.text)
        tf.close()
        tf = open(tmp_path, 'r', encoding='UTF-8')
        rusult = info_search(tf.read())
        tf.close()
    except Exception as e:
        print(e)
        rusult = {"username": [], "id": [], "ph_no": [],
                  "em": [], "address": [], "bir": []}
        pass
    try:
        # 刪除暫存檔
        # os.remove(pdf_path)
        os.remove(tmp_path)
        os.remove(pptx_path)
    except:
        pass
    return rusult


def read_txt(txt_path):
    try:
        tf = open(txt_path, 'r', encoding='UTF-8')
        rusult = info_search(tf.read())
        tf.close()
    except Exception as e:
        print(e)
        rusult = {"username": [], "id": [], "ph_no": [],
                  "em": [], "address": [], "bir": []}
        pass
    try:
        # 刪除暫存檔
        os.remove(txt_path)
    except:
        pass
    return rusult


# 直接執行測試

if __name__ == '__main__':
    print(read_xlsx('./tmp/tmp.xls'))
    print(read_xlsx('./tmp/tmp.xlsx'))
